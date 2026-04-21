import os
import json
import streamlit as st
from docx import Document
from pptx import Presentation
import PyPDF2
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Side, Font
from openpyxl.utils import get_column_letter
from openai import OpenAI
import datetime
import base64
import tempfile

# ===================== 1. 登录密码保护（第一道锁） =====================
def check_login():
    if "login_pass" not in st.session_state:
        st.session_state["login_pass"] = False

    if not st.session_state["login_pass"]:
        st.set_page_config(page_title="AI战略分析工具", layout="centered")
        st.title("🔒 AI战略分析工具 管理员登录")
        login_pwd = st.text_input("请输入登录密码", type="password")

        # 你自己设置的登录密码（可修改）
        YOUR_LOGIN_PASSWORD = "Ai@2026666"

        if st.button("登录", type="primary"):
            if login_pwd == YOUR_LOGIN_PASSWORD:
                st.session_state["login_pass"] = True
                st.rerun()
            else:
                st.error("密码错误，无权访问")
        return False
    return True


# 未登录直接拦截
if not check_login():
    st.stop()

# ===================== 2. API密钥配置（第二道锁） =====================
st.set_page_config(page_title="AI战略分析工具", layout="wide")
st.title("📊 AI战略分析表生成工具")

# 优先从Streamlit Secrets读取密钥（部署用），也支持手动输入（本地用）
# with st.sidebar:
#     st.header("⚙️ API配置")
#     # 部署时把密钥填在Streamlit Secrets里，这里自动读取，不用手动输
#     default_key = st.secrets.get("API_KEY", "")
#     api_key = st.text_input("通义千问API密钥", value=default_key, type="password")
#     if not api_key:
#         st.warning("请输入API密钥后使用")
#         st.stop()
api_key = st.secrets["API_KEY"]

BASE_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1"
MODEL_NAME = "qwen3-max"
VISION_MODEL = "qwen3-vl-flash"

# ===================== 3. 文件读取功能 =====================
# def read_file(file_bytes, filename):
#     ext = os.path.splitext(filename)[1].lower()
#     content = ""
#     try:
#         if ext == ".docx":
#             doc = Document(file_bytes)
#             content = "\n".join([p.text for p in doc.paragraphs])
#         elif ext == ".pptx":
#             prs = Presentation(file_bytes)
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         content += shape.text + "\n"
#         elif ext == ".pdf":
#             import pdfplumber
#             with pdfplumber.open(file_bytes) as pdf:
#                 for page in pdf.pages:
#                 # 按位置智能排序，解决分栏、错位问题
#                     text = page.extract_text(x_tolerance=2, y_tolerance=5)
#                     if text:
#                         content += text + "\n\n"
#         elif ext in [".xlsx", ".xls"]:
#             from openpyxl import load_workbook
#             wb = load_workbook(file_bytes, read_only=True)
#             for sheet in wb:
#                 for row in sheet.iter_rows(values_only=True):
#                     row_str = " ".join([str(cell) for cell in row if cell is not None])
#                     if row_str:
#                         content += row_str + "\n"
#     except Exception as e:
#         st.error(f"文件{filename}读取失败：{str(e)}")
#     return content
# ===================== 3. 文件读取功能（全AI提取版，替换原有函数） =====================
def read_file(file_bytes, filename):
    ext = os.path.splitext(filename)[1].lower()
    content = f"===== {filename} =====\n"

    # 写入临时文件
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(file_bytes.getbuffer())
        tmp_path = tmp.name

    try:
        client = OpenAI(api_key=api_key, base_url=BASE_URL)

        # PDF：视觉大模型识图提取（豆包同款）
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(tmp_path) as pdf:
                for idx, page in enumerate(pdf.pages, 1):
                    # 转图片
                    img = page.to_image()
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as img_tmp:
                        img.save(img_tmp.name)
                        with open(img_tmp.name, "rb") as f:
                            b64 = base64.b64encode(f.read()).decode()

                    # 通义视觉AI提取
                    messages = [
                        {
                            "role": "user",
                            "content": [
                                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                                {"type": "text", "text": "逐字提取所有文字，保持视觉顺序，不要修改"}
                            ]
                        }
                    ]
                    res = client.chat.completions.create(model=VISION_MODEL, messages=messages)
                    content += f"\n第{idx}页：{res.choices[0].message.content}\n"
                    os.remove(img_tmp.name)

        # Word/PPT/Excel：文本大模型标准化提取
        else:
            raw_content = ""
            if ext == ".docx":
                doc = Document(tmp_path)
                raw_content = "\n".join([p.text for p in doc.paragraphs])
            elif ext == ".pptx":
                prs = Presentation(tmp_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            raw_content += shape.text + "\n"
            elif ext in [".xlsx", ".xls"]:
                from openpyxl import load_workbook
                wb = load_workbook(tmp_path, read_only=True)
                for sheet in wb:
                    for row in sheet.iter_rows(values_only=True):
                        row_str = " ".join([str(c) for c in row if c is not None])
                        raw_content += row_str + "\n"

            # AI 整理内容
            prompt = f"完整提取文档内容，保持清晰排版：\n{raw_content}"
            res = client.chat.completions.create(model=MODEL_NAME, messages=[{"role": "user", "content": prompt}])
            content += res.choices[0].message.content

    except Exception as e:
        content += f"读取失败：{str(e)}"
    finally:
        os.remove(tmp_path)

    return content


# ===================== 【新增】4. 战略条目相关性判断核心函数 =====================
def get_correlation_matrix(row_items, col_items):
    """
    行条目：主要改进事项 + 战略目标2030（Excel行维度）
    列条目：年度目标2026 + 改进指标（Excel列维度）
    AI自动判断每对条目是否强相关，返回需要打勾的(行索引, 列索引)列表
    """
    # 无条目直接返回空，避免报错
    if not row_items or not col_items:
        return []
    
    # 整理AI输入内容，带索引标记
    row_texts = [f"【行{i}】{item['text']}" for i, item in enumerate(row_items)]
    col_texts = [f"【列{j}】{item['text']}" for j, item in enumerate(col_items)]
    
    # 精准prompt，严格约束输出格式
    system_prompt = """你是专业的战略规划分析师，精准判断行条目和列条目之间的业务相关性。
判断规则：
1. 仅当行条目和列条目存在**直接、强业务关联、因果支撑关系、落地对应关系**时，才判定为强相关。
2. 弱相关、间接关联、无关联的条目，一律不判定。
3. 严格返回JSON格式，结构固定为：{"correlations": [[行索引, 列索引], [行索引, 列索引], ...]}
4. 只返回符合强相关的索引对，禁止额外解释、多余内容。
"""
    user_prompt = f"""
所有行维度条目（主要改进事项+战略目标）：
{chr(10).join(row_texts)}

所有列维度条目（年度目标+改进指标）：
{chr(10).join(col_texts)}

请严格按照规则，返回强相关的行索引和列索引对。
"""
    try:
        client = OpenAI(api_key=api_key, base_url=BASE_URL)
        resp = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            response_format={"type": "json_object"},
            temperature=0.1  # 低温保证判断稳定一致
        )
        result = json.loads(resp.choices[0].message.content)
        return result.get("correlations", [])
    except Exception as e:
        st.error(f"相关性自动判断失败：{str(e)}，已跳过打勾步骤")
        return []

# ===================== 4. AI分析功能 =====================
def analyze_with_ai(content, feedback=None):
    system_prompt = """你是专业战略分析师，提取4类内容，其中战略目标3-6条，年度目标8-12条，主要改进事项15-20条，改进指标≥20条，用中文简洁短句，严格返回JSON：
{
    "战略目标2030": [],
    "年度目标2026": [],
    "主要改进事项": [],
    "改进指标": []
}"""
    messages = [{"role": "system", "content": system_prompt}]
    # 加载历史对话
    hist = st.session_state.get("current_session", {}).get("history", [])
    for h in hist:
        messages.append({"role": h["role"], "content": h["text"]})
    # 拼接当前输入
    # 有修改意见时：把【原始文档 + 修改要求】一起发给AI
    if feedback:
        combined_prompt = f"参考以下原始文档内容：\n{content}\n\n修改要求：{feedback}"
        messages.append({"role": "user", "content": combined_prompt})
    # 无修改意见时（首次生成）：只传文档内容
    else:
        messages.append({"role": "user", "content": content})
    client = OpenAI(api_key=api_key, base_url=BASE_URL)
    resp = client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        response_format={"type": "json_object"}
    )
    return json.loads(resp.choices[0].message.content)


# ===================== 5. Excel生成功能 =====================
def save_excel(data, base_name="分析结果"):
    out = f"{base_name}_战略表.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "战略规划"
    thin = Side(style='thin', color='000000')
    border = Border(top=thin, bottom=thin, left=thin, right=thin)
    font = Font(name="宋体", size=11)

    # 中心单元格
    CR, CC = 30, 30
    ws.cell(CR, CC, value="年度目标2026")

       # ========== 修改点1：填充四个方向时，记录条目坐标与内容，用于后续相关性判断 ==========
    row_items = []  # 行维度：主要改进事项(上) + 战略目标2030(下)
    col_items = []  # 列维度：年度目标2026(左) + 改进指标(右)

    # 上：主要改进事项（行维度，从上到下排列）
    items = data.get("主要改进事项", [])[:25]  # 与原代码截断逻辑保持一致
    r, c = CR - 1, CC
    for idx, x in enumerate(items):
        cell = ws.cell(r, c, x)
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        row_items.append({"text": x, "row": r})  # 记录行号与内容
        r -= 1

    # 下：战略目标2030（行维度，追加到行条目列表）
    items = data.get("战略目标2030", [])[:25]
    r, c = CR + 1, CC
    for idx, x in enumerate(items):
        cell = ws.cell(r, c, x)
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        row_items.append({"text": x, "row": r})  # 记录行号与内容
        r += 1

    # 左：年度目标2026（列维度，从右到左排列）
    items = data.get("年度目标2026", [])[:25]
    r, c = CR, CC - 1
    for idx, x in enumerate(items):
        cell = ws.cell(r, c, x)
        cell.alignment = Alignment(textRotation=90, horizontal='center', vertical='center')
        col_items.append({"text": x, "col": c})  # 记录列号与内容
        c -= 1

    # 右：改进指标（列维度，追加到列条目列表）
    items = data.get("改进指标", [])[:25]
    r, c = CR, CC + 1
    for idx, x in enumerate(items):
        cell = ws.cell(r, c, x)
        cell.alignment = Alignment(textRotation=90, horizontal='center', vertical='center')
        col_items.append({"text": x, "col": c})  # 记录列号与内容
        c += 1

    # ========== 修改点2：调用AI相关性判断，自动打勾 ==========
    if row_items and col_items:
        correlations = get_correlation_matrix(row_items, col_items)
        # 填充打勾单元格
        for row_idx, col_idx in correlations:
            # 索引合法性校验，避免越界报错
            if 0 <= row_idx < len(row_items) and 0 <= col_idx < len(col_items):
                target_row = row_items[row_idx]["row"]
                target_col = col_items[col_idx]["col"]
                check_cell = ws.cell(target_row, target_col, "√")
                # 打勾样式：红色加粗居中，更醒目
                check_cell.alignment = Alignment(horizontal='center', vertical='center')
                check_cell.font = Font(name="宋体", size=11, bold=True, color="FF0000")

    # 清理空行空列
    for r in range(ws.max_row, 0, -1):
        if all(ws.cell(r, c).value is None for c in range(1, ws.max_column + 1)):
            ws.delete_rows(r)
    for c in range(ws.max_column, 0, -1):
        if all(ws.cell(r, c).value is None for r in range(1, ws.max_row + 1)):
            ws.delete_cols(c)

    # 定位中心单元格
    ar, ac = CR, CC
    for r in range(1, ws.max_row + 1):
        for c in range(1, ws.max_column + 1):
            if ws.cell(r, c).value == "年度目标2026":
                ar, ac = r, c
                break

    # 插入中心图片four.png
    img_path = "four.png"
    if os.path.exists(img_path):
        from openpyxl.drawing.image import Image
        img = Image(img_path)
        img.width = 400
        img.height = 400
        col_letter = get_column_letter(ac)
        ws.add_image(img, f"{col_letter}{ar}")

    # 设置行高列宽
    for r in range(1, ws.max_row + 1):
        ws.row_dimensions[r].height = 15
    ws.row_dimensions[ar].height = 300
    for c in range(1, ws.max_column + 1):
        ws.column_dimensions[get_column_letter(c)].width = 3
    ws.column_dimensions[get_column_letter(ac)].width = 50

    # 统一设置所有单元格的边框和字体（包含打勾的单元格）
    for r in range(1, ws.max_row + 1):
        for c in range(1, ws.max_column + 1):
            cell = ws.cell(r, c)
            cell.border = border
            # 保留打勾单元格的特殊字体，其他单元格用默认字体
            if cell.value != "√":
                cell.font = font

    wb.save(out)
    return out


# ===================== 6. 主界面 =====================
# 初始化会话状态（新增存储原始文档内容）
if "current_session" not in st.session_state:
    st.session_state["current_session"] = {
        "history": [],
        "last_data": None,
        "original_content": ""  # 新增：永久保存上传的文档全文
    }

# 文件上传区
st.subheader("📁 上传文件（支持多选：Word/PPT/PDF/Excel）")
uploaded_files = st.file_uploader(
    "选择文件",
    type=["docx", "pptx", "pdf", "xlsx", "xls"],
    accept_multiple_files=True
)

# 生成按钮
if st.button("🚀 生成Excel", type="primary"):
    if not uploaded_files:
        st.warning("请先上传文件！")
    else:
        with st.spinner("AI正在分析文档..."):
            # 合并所有文件内容
            all_content = ""
            for f in uploaded_files:
                all_content += f"\n===== 文件：{f.name} =====\n"
                all_content += read_file(f, f.name)
            # AI分析
            # 保存原始文档全文（只在第一次生成时存储）
            st.session_state["current_session"]["original_content"] = all_content
            # AI分析
            data = analyze_with_ai(all_content)
            st.session_state["current_session"]["last_data"] = data
            # 生成Excel
            out_path = save_excel(data)
            # 下载按钮
            with open(out_path, "rb") as f:
                st.download_button(
                    label="📥 下载生成的Excel",
                    data=f,
                    file_name=out_path,
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

# 修改意见区
st.subheader("✍️ 修改意见（重新生成）")
feedback = st.text_area("输入你的修改要求，AI会基于上一版调整")
if st.button("发送并重新生成Excel"):
    if not feedback:
        st.warning("请输入修改意见！")
    elif not st.session_state["current_session"]["last_data"]:
        st.warning("请先生成一次Excel！")
    else:
        with st.spinner("AI正在修改并重新生成..."):
            # 记录历史对话
            st.session_state["current_session"]["history"].append({"role": "user", "text": feedback})
            # AI重新生成
            # 读取保存的原始文档，结合修改意见一起发给AI
            original_content = st.session_state["current_session"]["original_content"]
            data = analyze_with_ai(original_content, feedback=feedback)
            st.session_state["current_session"]["last_data"] = data
            # 生成新Excel
            out_path = save_excel(data)
            # 下载按钮
            with open(out_path, "rb") as f:
                st.download_button(
                    label="📥 下载修改后的Excel",
                    data=f,
                    file_name=out_path,
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

# 历史对话区
with st.expander("📜 查看历史对话记录"):
    for item in st.session_state["current_session"]["history"]:
        st.write(f"**你**：{item['text']}")
